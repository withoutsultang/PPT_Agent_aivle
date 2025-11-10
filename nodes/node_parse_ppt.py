import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from utils.text_utils import clean_text

def node_parse_ppt(state):
    ppt = Presentation(state['pptx_path'])
    slide_idx = state.get("slide_index", 0)
    work_dir = state.get("work_dir", "./")
    texts, tables, images = [], [], []
    sd0 = ppt.slides[slide_idx]

    for i, sh in enumerate(sd0.shapes):
        if sh.has_text_frame:
            txt = "\n".join(p.text for p in sh.text_frame.paragraphs)
            texts.append(clean_text(txt))
        if sh.shape_type == MSO_SHAPE_TYPE.TABLE:
            tbl = [[clean_text(c.text) for c in r.cells] for r in sh.table.rows]
            tables.append(tbl)
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            ext = sh.image.ext
            path = os.path.join(work_dir, "media", f"slide{slide_idx}_img_{i}.{ext}")
            os.makedirs(os.path.dirname(path), exist_ok=True)
            with open(path, "wb") as f:
                f.write(sh.image.blob)
            images.append(path)

    state["texts"] = texts
    state["tables"] = tables
    state["images"] = images
    return state
