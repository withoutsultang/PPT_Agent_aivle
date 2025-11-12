# agent_nodes.py

import os, re, textwrap, subprocess, json, time
from typing import List, Dict, Optional, TypedDict, Any
from openai import OpenAI
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

# 로컬 모듈 임포트
from utils import clean_text, split_sents, ffprobe_duration, img_to_data_url, render_mp4, concat_videos_ffmpeg, export_slide_as_png

# --- 환경 설정 ---
LLM_MODEL = "gpt-4o-mini"
TTS_MODEL = "tts-1" # TTS-1-HD가 더 고음질이나, tts-1이 더 빠르고 비용 효율적
client = OpenAI()

# --- State 정의 ---
class State(TypedDict, total=False):
  pptx_path: str
  work_dir: str
  prompt: Dict
  slide_index: int
  total_slides: int # 추가: 총 슬라이드 수

  titles: List[str]
  texts: List[str]
  tables: List[List[List[str]]]
  images: List[str]
  shape_texts: List[str] # 도형 텍스트를 저장할 필드 추가
  slide_image: List[str]
  
  external_content: Dict[str, List[Dict[str, str]]]

  page_content: str
  script: str
  all_scripts: List[str] # 누적 스크립트
  quiz_set: List[Dict[str, Any]]
  
  audio: str
  video_path: List[str]
  video_paths: List[str]
  final_video: str
  
  failed_slides: List[int] # 실패한 슬라이드 인덱스 저장

# ===============================
# 🔹 Node Functions
# ===============================

def get_shapes_text(shape):
    """하나의 도형(또는 그룹)에서 텍스트를 재귀적으로 추출"""
    texts = []
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sh in shape.shapes:
            texts.extend(get_shapes_text(sh))
    elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and shape.has_text_frame:
        text = shape.text.strip()
        if text:
            texts.append(text)
    return texts

def node_parse_all(state: State) -> State:
    """PPT 파일에서 모든 슬라이드 정보를 추출하고 이미지로 변환 (1회 실행)"""
    
    ppt = Presentation(state['pptx_path'])
    work_dir = state.get("work_dir", "./")

    # 미디어/슬라이드 이미지 저장 경로 설정
    MEDIA_DIR = os.path.join(work_dir, "media")
    SLIDES_DIR = os.path.join(work_dir, "slides")
    os.makedirs(MEDIA_DIR, exist_ok=True)
    os.makedirs(SLIDES_DIR, exist_ok=True)

    texts, tables, images, titles, slide_image, shapes = [], [], [], [], [], []

    for slide_idx, slide in enumerate(ppt.slides):
        # 1. 슬라이드 이미지(스냅샷) 추출
        slide_state = {"pptx_path": state['pptx_path'], "work_dir": SLIDES_DIR, "slide_index": slide_idx}
        slide_state = export_slide_as_png(slide_state)
        
        src_path = slide_state["slide_image"]
        dst_path = os.path.join(SLIDES_DIR, f"slide_img{slide_idx+1}.png")
        if os.path.exists(src_path):
            os.replace(src_path, dst_path) # 파일 이동
            slide_image.append(dst_path)
        else:
            slide_image.append(None)

        # 2. 텍스트, 표, 이미지 정보 추출
        full_slide_text, slide_tables, slide_images, slide_title, slide_shapes_texts = "", [], [], "", []
        
        for sh in slide.shapes:
            if sh.is_placeholder and sh.placeholder_format.type == PP_PLACEHOLDER.TITLE and sh.has_text_frame:
                slide_title = sh.text.strip()

            if sh.has_text_frame:
                full_slide_text += "\\n".join(p.text for p in sh.text_frame.paragraphs) + "\\n"
            
            if sh.shape_type == MSO_SHAPE_TYPE.TABLE:
                tbl = [[clean_text(c.text) for c in r.cells] for r in sh.table.rows]
                slide_tables.append(tbl)
            
            # 도형 텍스트 추출 (재귀)
            slide_shapes_texts.extend(get_shapes_text(sh))

            if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                ext = sh.image.ext
                img_filename = f"slide{slide_idx+1}_img_{len(slide_images)}.{ext}"
                path = os.path.join(MEDIA_DIR, img_filename)
                slide_images.append(path)
                with open(path, "wb") as f:
                    f.write(sh.image.blob)

        # 3. 결과 누적
        texts.append(clean_text(full_slide_text))
        tables.append(slide_tables)
        images.append(slide_images)
        titles.append(slide_title)
        shapes.append(",".join(slide_shapes_texts))

    # 4. State 저장
    state.update({
        'texts': texts,
        'tables': tables,
        'images': images,
        'slide_image': slide_image,
        'titles': titles,
        'shape_texts': shapes,
        "total_slides": len(ppt.slides)
    })
    
    return state

def serpapi_search_by_title(title: str, num: int = 4) -> list[dict]:
    """SerpAPI를 이용해 실제 검색을 수행하고 필터링된 결과를 반환"""
    key = os.getenv("SERPAPI_API_KEY")
    EXCLUDE_DOMAINS = ["blog.naver.com", "tistory.com", "brunch.co.kr", "medium.com", "velog.io", "kin.naver.com", "reddit.com", "youtube.com"]
    query = f"{title} " + " ".join([f"-site:{d}" for d in EXCLUDE_DOMAINS])
    
    try:
        res = requests.get("https://serpapi.com/search.json", params={
            "engine": "google", "q": query, "hl": "ko", "gl": "kr", "num": num, "api_key": key
        }, timeout=15)
        
        data = res.json().get("organic_results", []) or []
        results = []
        for item in data:
            url = item.get("link", "")
            if not url: continue
            domain = urlparse(url).netloc
            results.append({
                "title": item.get("title", ""),
                "url": url,
                "snippet": item.get("snippet", ""),
                "domain": domain
            })
        return results
    except Exception as e:
        print(f"[SerpAPI 오류] 검색 실패: {e}")
        return []

def node_tool_search(state: dict) -> dict:
    """외부 검색 노드: 슬라이드 제목을 기반으로 검색을 수행하고 결과를 state에 저장"""
    idx = state.get("slide_index", 0)
    titles = state.get("titles", [])
    texts_all = state.get("texts", [])
    tables_all = state.get("tables", [])
    images_all = state.get("images", [])

    title = titles[idx] if idx < len(titles) else ""
    texts = texts_all[idx] if idx < len(texts_all) else ""
    
    state["external_content"] = {"queries": [], "summaries": [], "references": []} # 초기화

    # 쿼리 생성 로직 (생략: 기존 코드와 동일)
    queries = []
    if title: queries.append({"text": title, "context": "title"})
    if title and texts: queries.append({"text": f"{title} {texts[:80]}", "context": "title+text"})
    # ... (필요에 따라 table, image 쿼리 추가 로직)
    
    # 검색 수행
    all_results = []
    for q in queries:
        results = serpapi_search_by_title(q["text"], num=4)
        all_results.extend(results)
        time.sleep(0.2)
        
    # 결과 정리 (중복 제거 및 구조화)
    summaries = [{"text": clean_text(r["snippet"]), "source": r["title"]} for r in all_results if r.get("snippet")]
    references = [{"title": clean_text(r["title"]), "url": r["url"]} for r in all_results]
        
    state["external_content"] = {
        "queries": queries,
        "summaries": summaries[:3], # 상위 3개만 요약에 사용
        "references": references[:4] # 상위 4개만 참조 출처로 사용
    }
    return state

def node_generate_page_content(state: State) -> State:
    """LLM을 호출하여 현재 슬라이드 정보와 외부 자료를 통합하여 페이지 설명문 생성"""
    idx        = int(state.get("slide_index", 0))
    titles     = state.get("titles", [])
    texts_all  = state.get("texts", [])
    tables_all = state.get("tables", [])
    images_all = state.get("images", [])
    shapes_all = state.get("shape_texts", [])
    prompt     = clean_text(state.get("prompt", {}).get("style", "")) # style 프롬프트

    title  = clean_text(str(titles[idx])) if idx < len(titles) else ""
    texts  = clean_text(str(texts_all[idx])) if idx < len(texts_all) else ""
    tables = tables_all[idx] if idx < len(tables_all) else []
    images = images_all[idx] if idx < len(images_all) else []
    shapes = shapes_all[idx] if idx < len(shapes_all) else ""

    # 표 전처리
    table_text = ""
    if tables and tables[0]:
        first_table = tables[0][:6] 
        table_text = "\\n".join([f"| {' | '.join(row)} |" for row in first_table])

    # 이미지 인코딩 (최대 3장)
    image_data_urls = [img_to_data_url(p) for p in images[:3] if os.path.exists(p)]

    # 외부 보완 블록 구성
    ext = state.get("external_content", {}) or {}
    ext_refs      = ext.get("references", [])
    ext_summaries = ext.get("summaries", [])
    
    ext_ref_block = "\\n".join([f"[{i+1}] {r.get('title','')} — {r.get('url','')}" for i, r in enumerate(ext_refs)])
    ext_summary_block = "\\n".join([f"- {s.get('text')} ({s.get('source')})" for s in ext_summaries])
    
    # 프롬프트 구성
    content_input = textwrap.dedent(f"""
        다음은 한 슬라이드의 정보와 외부 보완 자료입니다.
        제목: {title}
        ---
        [텍스트]: {texts}
        [표]:\\n{table_text}
        [도형/객체 텍스트 요약]: {shapes}
        [프롬프트 스타일 지침]: {prompt}
        ---
        [외부 보완 자료]
        - 핵심 보강 요약:\\n{ext_summary_block}
        - 참조 출처:\\n{ext_ref_block}
        ---
        규칙:
        1) 모든 정보를 통합하여 **4~6문장**의 간결한 **슬라이드 설명문**을 작성할 것.
        2) 표/이미지/도형 의미를 자연스럽게 통합해 설명.
        3) 외부 보완 내용은 핵심만 반영하며, 출처를 대괄호 숫자로 표시 (예: [1][2]).
    """)

    # LLM 호출
    messages = [{"role": "system", "content": "당신은 슬라이드의 모든 정보를 통합하여 핵심 내용을 요약하는 전문 에이전트입니다."},
                {"role": "user", "content": [{"type": "text", "text": content_input}]}]
    for img_url in image_data_urls:
        messages[-1]["content"].append({"type": "image_url", "image_url": {"url": img_url}})

    response = client.chat.completions.create(model=LLM_MODEL, messages=messages, temperature=0.5)

    # 결과 저장
    page_content = clean_text(response.choices[0].message.content)
    state["page_content"] = " ".join(split_sents(page_content))
    return state

def node_generate_script(state: State) -> State:
    """강의 스크립트 생성: 이전 스크립트와 다음 목차를 고려하여 연속성 있게 작성"""
    
    all_titles = state.get("titles", [])
    prev_scripts = state.get("all_scripts", [])
    previous_script = prev_scripts[-1] if prev_scripts else "없음"
    
    prompt_data = state.get("prompt", {})
    tone = prompt_data.get("tone", "친절하고 명료한 강의 톤")
    target_time = prompt_data.get("target_duration_sec", 60)
    current_page_content = state.get("page_content", "")

    current_index = state.get("slide_index", 0)
    total_slides = state.get("total_slides", len(all_titles))
    current_title = all_titles[current_index] if all_titles and current_index < len(all_titles) else "현재 슬라이드"
    
    # --- 강의 흐름(Flow) 지시사항 구성 (가장 중요한 고도화 파트) ---
    flow_instruction = ""
    if current_index == 0:
        flow_instruction = "이것은 전체 강의의 '첫 번째' 슬라이드입니다. 강의 전체를 소개하는 '도입부'로 시작하되, 청중에게 직접적인 인사말은 넣지 말아 주세요. 하나의 긴 강의가 시작되는 것처럼 자연스럽게 시작해야 합니다. (이전 내용 요약/마무리 인사는 금지)"
    elif current_index == total_slides - 1:
        flow_instruction = "이것은 전체 강의의 '마지막' 슬라이드입니다. 강의 전체를 요약하고 청중에게 '마무리 끝인사'를 반드시 포함해 주세요. (다음 내용 예고 금지)"
    else:
        next_title = all_titles[current_index + 1] # 다음 슬라이드의 제목을 가져옴
        last_sentence_part = previous_script[-50:] if previous_script != "없음" and len(previous_script) > 50 else previous_script
        
        flow_instruction = (
            f"이것은 강의의 '중간' 슬라이드입니다. 직전 슬라이드 스크립트의 마지막 문장(예: '...{last_sentence_part}')에서 내용이 '완벽하게 연결'되도록 현재 슬라이드의 설명을 바로 시작해 주세요. "
            f"스크립트의 마지막 부분에 다음 슬라이드의 주제인 '[{next_title}]'를 활용하여 청중의 기대감을 높이는 자연스러운 연결 및 예고 멘트를 포함해야 합니다. "
            "별도의 연결 멘트 없이 바로 본론을 시작하며, 하나의 긴 강의처럼 흐름을 유지해야 합니다. (이전 내용 요약은 금지, 다음 내용 예고는 필수)"
        )

    # --- LLM 프롬프트 설계 ---
    system_prompt = (
        "당신은 전문 AI 강사입니다. 이 스크립트는 여러 슬라이드를 연결하여 제작될 '하나의 연속적인 강의 영상'에 사용될 것입니다. "
        "강의의 전체 목차와 흐름을 고려하여, 모든 슬라이드 스크립트가 끊김 없이 매끄럽게 이어지도록 작성해야 합니다."
    )

    user_prompt = f"""
    # 전체 강의 목차
    {all_titles}

    # 현재 강의 중인 슬라이드
    - 인덱스: {current_index}
    - 제목: {current_title}

    # 직전 슬라이드 스크립트
    {previous_script}

    #현재 슬라이드 핵심 내용
    {current_page_content}

    # 필수) 스크립트 작성 조건
    1. 톤앤매너: {tone}
    2. 분량: 약 {target_time}초 분량의 스크립트를 작성해 주세요. (TTS 재생 속도 1.0x 기준)
    3. [중요] 흐름 지시사항: {flow_instruction}
    4. [연속성 규칙] '오늘', '이번 강의에서는', '안녕하세요', '마지막으로', '감사합니다' 등 강의의 연속성을 끊거나 시간/날짜를 특정하는 표현은 마지막 슬라이드의 최종 끝인사를 제외하고는 **절대 사용하지 마세요.**
    5. [생동감] 청중의 이해를 돕기 위해 현재 슬라이드의 내용 중 중요한 부분이나 그래프/이미지를 언급하며 '청중에게 말을 거는 듯한' 구어체와 생동감을 불어 넣어주세요.
    6. [근거 제시] 슬라이드에 제시된 데이터(그래프, 표, 수치)나 검색된 외부 정보(예: Amazon SageMaker)를 언급할 때는 "화면의 그래프에서", "이 표에서 확인하실 수 있듯이", "Amazon SageMaker와 같은 플랫폼을 예로 들면" 등의 표현으로 근거를 제시하며 설명해 주세요.

    [스크립트 시작]
    """

    # (3) LLM 호출
    response = client.chat.completions.create(
        model=LLM_MODEL, messages=[{"role": "system", "content": system_prompt}, {"role": "user", "content": user_prompt}], temperature=0.7
    )

    script = clean_text(response.choices[0].message.content).replace("[스크립트 시작]", "").replace("[스크립트 종료]", "")
    
    # State 업데이트
    state["script"] = script
    if "all_scripts" not in state: state["all_scripts"] = []
    state["all_scripts"].append(script)

    return state

def node_tts(state: dict) -> dict:
    """발표 스크립트를 음성(mp3)으로 변환하고 속도 조절"""
    script = state.get("script", "")
    prompt = state.get("prompt", {})
    voice = prompt.get("voice", "alloy").split('-')[-1].strip()
    work_dir = state.get("work_dir", "./")
    speed = float(prompt.get("speed", 1.0))
    slide_idx = int(state.get("slide_index", 0))

    if not script.strip(): raise ValueError("스크립트가 비어 있습니다.")

    os.makedirs(work_dir, exist_ok=True)
    base_audio_path = os.path.join(work_dir, f"narration_raw_{slide_idx}.mp3")
    final_audio_path = os.path.join(work_dir, f"narration_{slide_idx}_{speed}x.mp3")

    # OpenAI TTS 호출
    response = client.audio.speech.create(model=TTS_MODEL, voice=voice, input=script, response_format="mp3")
    with open(base_audio_path, "wb") as f:
        f.write(response.read())

    # FFmpeg로 속도 조절
    if speed != 1.0:
        # ffmpeg atempo 필터는 0.5x ~ 100.0x만 지원 (0.5 이하/2.0 초과는 atempo 필터 체인 필요)
        if speed > 2.0 or speed < 0.5: 
            # 0.5~2.0 범위를 벗어나면, 여러 atempo를 체인으로 연결하여 처리
            current_speed = speed
            atempo_filters = []
            while current_speed > 2.0:
                atempo_filters.append("atempo=2.0")
                current_speed /= 2.0
            while current_speed < 0.5:
                atempo_filters.append("atempo=0.5")
                current_speed /= 0.5
            if current_speed != 1.0:
                 atempo_filters.append(f"atempo={current_speed}")
            filter_chain = ",".join(atempo_filters)
            
        else:
            filter_chain = f"atempo={speed}"

        cmd = ["ffmpeg", "-y", "-i", base_audio_path, "-filter:a", filter_chain, "-b:a", "192k", final_audio_path]
        subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, check=True)
        
    else:
        final_audio_path = base_audio_path

    state["audio"] = final_audio_path
    
    return state

def node_make_video(state: dict) -> dict:
    """슬라이드 이미지와 음성을 합쳐 슬라이드별 MP4 영상 생성"""
    slide_imgs = state.get("slide_image", [])
    audio_path = state.get("audio", "")
    work_dir = state.get("work_dir", "./")
    slide_index = state.get("slide_index", 0)

    if not slide_imgs or slide_index >= len(slide_imgs) or not os.path.exists(audio_path):
        return state

    if "video_path" not in state: state["video_path"] = []

    video_filename = f"slide{slide_index+1}_lecture.mp4"
    out_mp4 = os.path.join(work_dir, video_filename)

    # 실제 영상 생성
    render_mp4(image_path=slide_imgs[slide_index], audio_path=audio_path, out_mp4=out_mp4)
    
    # 중복 방지하여 video_path에 추가
    if out_mp4 not in state["video_path"]:
        state["video_path"].append(out_mp4)

    return state

def node_accumulate_and_step(state: dict) -> dict:
    """영상 누적 및 다음 슬라이드 인덱스 증가"""
    current_idx = state.get("slide_index", 0)
    total = state.get("total_slides", 1)
    video_list = state.get("video_path", [])

    # 누적 리스트 초기화
    if "video_paths" not in state or not isinstance(state["video_paths"], list):
        state["video_paths"] = []

    # 1️⃣ 영상 검증 및 누적
    if len(video_list) > current_idx:
        current_video = video_list[current_idx]
        if os.path.exists(current_video):
            if current_video not in state["video_paths"]:
                state["video_paths"].append(current_video)
        else:
            if "failed_slides" not in state: state["failed_slides"] = []
            state["failed_slides"].append(current_idx + 1) # 1-based index
    
    # 2️⃣ 다음 슬라이드로 이동
    state["slide_index"] = current_idx + 1

    return state

def router_continue_or_done(state: dict) -> str:
    """다음 슬라이드 유무에 따라 CONTINUE 또는 DONE 분기"""
    current_idx = state.get("slide_index", 0)
    total_slides = state.get("total_slides", len(state.get("titles", [])))

    if current_idx >= total_slides:
        return "done"
    else:
        return "continue"

def node_concat(state: State) -> State:
    """video_paths의 모든 영상을 순서대로 연결하여 최종 영상 생성"""
    video_paths = state.get("video_paths", [])
    work_dir = state.get("work_dir", "./step1_output")

    if not video_paths:
        return state

    final_video = os.path.join(work_dir, "final_lecture.mp4")
    
    # FFmpeg로 영상 합치기 (reencode=False로 빠르고 단순 복사)
    concat_videos_ffmpeg(video_paths=video_paths, out_path=final_video, reencode=False)

    state["final_video"] = final_video

    return state

def node_generate_quiz(state: dict) -> dict:
    """강의 스크립트 전체를 바탕으로 복습 퀴즈를 JSON 형식으로 생성"""
    all_scripts = state.get("all_scripts", [])

    if not all_scripts:
        state["quiz_set"] = []
        return state

    system_prompt = textwrap.dedent("""
        당신은 강의 내용을 복습시키는 전문 교육 보조입니다. 제공된 강의 스크립트 전체 내용을 바탕으로,
        핵심 내용을 확인할 수 있는 퀴즈 세트를 생성해야 합니다.
        퀴즈는 반드시 객관식 4지선다형이어야 하며, 강의의 핵심 개념을 다루어야 합니다.
        **반드시 유효한 JSON 형식으로만 응답해야 합니다.**
    """)

    full_lecture_script = "\\n\\n".join([f"[슬라이드 {i+1}]\\n{script}" for i, script in enumerate(all_scripts)])

    user_prompt = textwrap.dedent(f"""
        --- [강의 전체 내용] ---
        {full_lecture_script}
        ---

        [규칙]
        1. 위 강의 내용 전체를 바탕으로 **총 6개의 [객관식 퀴즈]**를 생성하세요.
        2. 각 퀴즈는 반드시 4개의 선택지(options)를 가져야 합니다.
        3. **[중요]** 각 선택지는 **'1. 선택지 내용', '2. 선택지 내용'** 처럼 반드시 번호로 시작해야 합니다.
        4. 각 퀴즈마다 [question], [options], [answer] 키만 포함해야 합니다.
        5. **[중요]** 정답(answer)은 **번호가 포함된 선택지 텍스트와 정확히 일치**해야 합니다. (예: "1. 선택지 1")
        6. 출력은 반드시 JSON 배열 형식이어야 합니다.

        [JSON 출력]
    """)

    try:
        response = client.chat.completions.create(
            model=LLM_MODEL, messages=[{"role": "system", "content": system_prompt}, {"role": "user", "content": user_prompt}],
            response_format={"type": "json_object"}
        )
        quiz_set_data = json.loads(response.choices[0].message.content.strip())
        # LLM이 직접 배열을 반환하거나, "quizzes" 등의 키로 감쌀 수 있음.
        state["quiz_set"] = quiz_set_data.get("quizzes", quiz_set_data) 
    except Exception as e:
        state["quiz_set"] = [] 
        print(f"[오류] 퀴즈 생성 또는 JSON 파싱 실패: {e}")
        
    return state